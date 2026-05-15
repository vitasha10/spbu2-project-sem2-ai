"""Генерация графиков и сборка презентации .pptx из шаблона СПбГУ.

Использует пустой шаблон base_empty.pptx из соседнего проекта - там сохранён
дизайн СПбГУ (титульный, рабочие, закрывающий), но без примеров слайдов.
"""
import sys
from pathlib import Path

import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
import seaborn as sns
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor

SPBGU_RED = RGBColor(0x95, 0x37, 0x34)

import joblib
from sklearn.calibration import CalibrationDisplay
from sklearn.inspection import permutation_importance
from sklearn.metrics import roc_auc_score

ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(ROOT))

from src.data import load_raw, clean
from src.features import add_features

TEMPLATE = 'C:/it/test-faceswap-2026/presentation/output/base_empty.pptx'
OUT_DIR = ROOT / 'docs'
IMG_DIR = OUT_DIR / 'pptx_assets'
OUT_PPTX = OUT_DIR / 'presentation.pptx'

sns.set_theme(style='whitegrid', font_scale=0.9)
plt.rcParams['axes.titlesize'] = 12
plt.rcParams['axes.labelsize'] = 10
plt.rcParams['xtick.labelsize'] = 9
plt.rcParams['ytick.labelsize'] = 9
plt.rcParams['font.family'] = 'DejaVu Sans'

MONTHS_RU = ['янв', 'фев', 'мар', 'апр', 'май', 'июн',
             'июл', 'авг', 'сен', 'окт', 'ноя', 'дек']


def save_charts():
    IMG_DIR.mkdir(parents=True, exist_ok=True)
    df_raw = load_raw(str(ROOT / 'data' / 'hotels.csv'))
    df = clean(df_raw)
    df = add_features(df)

    # 1. баланс классов
    fig, ax = plt.subplots(figsize=(5, 3.2), dpi=160)
    counts = df_raw['is_canceled'].value_counts().rename({0: 'Доехали', 1: 'Отменили'})
    counts.plot(kind='bar', ax=ax, color=['#4c72b0', '#dd8452'])
    ax.set_title('Сколько броней отменяют')
    ax.set_ylabel('кол-во броней')
    ax.set_xlabel('')
    ax.tick_params(axis='x', rotation=0)
    for p in ax.patches:
        ax.annotate(f'{int(p.get_height()):,}'.replace(',', ' '),
                    (p.get_x() + p.get_width() / 2, p.get_height()),
                    ha='center', va='bottom', fontsize=9)
    fig.tight_layout()
    fig.savefig(IMG_DIR / 'balance.png')
    plt.close(fig)

    # 2. срок до заезда: доля отмен по корзинам срока
    fig, ax = plt.subplots(figsize=(5.5, 3.2), dpi=160)
    bins = [0, 7, 30, 60, 120, 200, 400, 800]
    labels = ['0–7', '7–30', '30–60', '60–120', '120–200', '200–400', '400+']
    tmp = df_raw.copy()
    tmp['bucket'] = pd.cut(tmp['lead_time'], bins=bins, labels=labels, include_lowest=True)
    rates = tmp.groupby('bucket', observed=True)['is_canceled'].mean()
    bars = ax.bar(range(len(labels)), rates.values,
                  color=['#4c72b0', '#5b8db8', '#6b9bba', '#c98564',
                         '#dd8452', '#d06a3a', '#c44e52'])
    ax.set_xticks(range(len(labels)))
    ax.set_xticklabels(labels)
    ax.set_ylabel('доля отмен')
    ax.set_xlabel('дней между бронированием и заездом')
    ax.set_title('Чем больше срок до заезда, тем чаще отменяют')
    ax.set_ylim(0, max(rates.values) * 1.15)
    for bar, v in zip(bars, rates.values):
        ax.text(bar.get_x() + bar.get_width() / 2, v + 0.01,
                f'{v:.0%}', ha='center', va='bottom', fontsize=9)
    fig.tight_layout()
    fig.savefig(IMG_DIR / 'lead_time.png')
    plt.close(fig)

    # 3. депозит
    fig, ax = plt.subplots(figsize=(5, 3.2), dpi=160)
    name_map = {'No Deposit': 'Без депозита', 'Refundable': 'Возвратный',
                'Non Refund': 'Невозвратный'}
    g = df_raw.groupby('deposit_type')['is_canceled'].mean()
    g.index = [name_map.get(x, x) for x in g.index]
    g = g.sort_values()
    g.plot(kind='barh', ax=ax, color='#c44e52')
    ax.set_xlabel('доля отмен')
    ax.set_ylabel('')
    ax.set_title('Доля отмен по типу депозита')
    for i, v in enumerate(g.values):
        ax.text(v + 0.01, i, f'{v:.0%}', va='center', fontsize=9)
    ax.set_xlim(0, 1.1)
    fig.tight_layout()
    fig.savefig(IMG_DIR / 'deposit.png')
    plt.close(fig)

    # 4. топ стран с русскими подписями
    country_ru = {
        'PRT': 'Португалия', 'GBR': 'Великобритания', 'FRA': 'Франция',
        'ESP': 'Испания', 'DEU': 'Германия', 'ITA': 'Италия',
        'IRL': 'Ирландия', 'BEL': 'Бельгия', 'BRA': 'Бразилия',
        'NLD': 'Нидерланды', 'USA': 'США', 'CHE': 'Швейцария',
    }
    fig, ax = plt.subplots(figsize=(6.2, 3.2), dpi=160)
    top = df_raw['country'].value_counts().head(10).index
    sub = df_raw[df_raw['country'].isin(top)]
    g = sub.groupby('country')['is_canceled'].mean().loc[top].sort_values()
    g.index = [f'{c} — {country_ru.get(c, c)}' for c in g.index]
    g.plot(kind='barh', ax=ax, color='#4c72b0')
    ax.set_xlabel('доля отмен')
    ax.set_ylabel('')
    ax.set_title('Доля отмен по странам клиента (топ-10)')
    for i, v in enumerate(g.values):
        ax.text(v + 0.005, i, f'{v:.0%}', va='center', fontsize=9)
    ax.set_xlim(0, max(g.values) * 1.18)
    fig.tight_layout()
    fig.savefig(IMG_DIR / 'country.png')
    plt.close(fig)

    # 5. сезонность
    fig, ax = plt.subplots(figsize=(5, 3.2), dpi=160)
    months_order = ['January', 'February', 'March', 'April', 'May', 'June',
                    'July', 'August', 'September', 'October', 'November', 'December']
    by_month = df_raw.groupby('arrival_date_month')['is_canceled'].mean().reindex(months_order)
    by_month.plot(kind='line', marker='o', ax=ax, color='#c44e52')
    ax.set_ylabel('доля отмен')
    ax.set_xlabel('')
    ax.set_title('Отмены по месяцам заезда')
    ax.set_xticks(range(12))
    ax.set_xticklabels(MONTHS_RU, rotation=0)
    fig.tight_layout()
    fig.savefig(IMG_DIR / 'season.png')
    plt.close(fig)

    # 6. сравнение моделей (горизонтальный bar для читаемости подписей)
    fig, ax = plt.subplots(figsize=(6.5, 3.4), dpi=160)
    models = ['Случайный ответ', 'Логистическая регрессия',
              'k ближайших соседей', 'Решающее дерево',
              'Случайный лес', 'Градиентный бустинг']
    aucs = [0.500, 0.840, 0.851, 0.881, 0.913, 0.920]
    colors = ['#999999', '#4c72b0', '#55a868', '#8172b3', '#dd8452', '#c44e52']
    bars = ax.barh(models, aucs, color=colors)
    ax.set_xlim(0.4, 1.0)
    ax.set_xlabel('качество ROC-AUC')
    ax.set_title('Сравнение моделей (5-фолдовая кросс-валидация)')
    ax.axvline(0.5, color='gray', ls='--', lw=0.8)
    for bar, v in zip(bars, aucs):
        ax.text(v + 0.005, bar.get_y() + bar.get_height() / 2,
                f'{v:.3f}', va='center', fontsize=9)
    ax.invert_yaxis()
    fig.tight_layout()
    fig.savefig(IMG_DIR / 'models.png')
    plt.close(fig)

    # 7. значимость признаков
    split = joblib.load(ROOT / 'models' / '_split.pkl')
    X_test = split['X_test']
    y_test = split['y_test']
    final = joblib.load(ROOT / 'models' / 'final_model.pkl')
    hgb = final['pipeline']

    rng = np.random.default_rng(42)
    sub_idx = rng.choice(len(X_test), 3000, replace=False)
    Xs = X_test.iloc[sub_idx]
    ys = y_test.iloc[sub_idx]
    imp = permutation_importance(hgb, Xs, ys, scoring='roc_auc',
                                 n_repeats=3, random_state=42, n_jobs=-1)
    importance_ru = {
        'lead_time': 'срок до заезда',
        'country': 'страна клиента',
        'deposit_type': 'тип депозита',
        'total_of_special_requests': 'число спец-запросов',
        'market_segment': 'канал продажи',
        'adr': 'цена за ночь',
        'arrival_date_year': 'год заезда',
        'arrival_date_week_number': 'неделя заезда',
        'arrival_date_day_of_month': 'день заезда',
        'previous_cancellations': 'прошлые отмены клиента',
        'agent': 'агентство',
        'customer_type': 'тип клиента',
        'required_car_parking_spaces': 'запрос парковки',
        'booking_changes': 'число изменений брони',
        'room_changed': 'сменили номер при заезде',
        'assigned_room_type': 'выданный номер',
        'reserved_room_type': 'запрошенный номер',
        'distribution_channel': 'канал дистрибуции',
        'adr_per_person': 'цена на гостя',
    }
    ser = pd.Series(imp.importances_mean, index=X_test.columns).sort_values(ascending=False).head(10)
    ser.index = [importance_ru.get(x, x) for x in ser.index]
    fig, ax = plt.subplots(figsize=(5.6, 3.4), dpi=160)
    ser.iloc[::-1].plot(kind='barh', ax=ax, color='#4c72b0')
    ax.set_xlabel('падение качества ROC-AUC\nпри перетасовке признака')
    ax.set_ylabel('')
    ax.set_title('Значимость признаков (топ-10)')
    fig.tight_layout()
    fig.savefig(IMG_DIR / 'importance.png')
    plt.close(fig)

    # 8. калибровка как residual-plot.
    # Ось Y — отклонение «реально_отменили − предсказали».
    # 0 = идеальная калибровка, > 0 = модель занижает, < 0 = задирает уверенность.
    # На таком графике реальные промахи в 5–8 % видны как заметные горбы,
    # а не сливаются с диагональю.
    rf = joblib.load(ROOT / 'models' / '_rf.pkl')
    lr = joblib.load(ROOT / 'models' / '_logreg.pkl')
    from sklearn.calibration import calibration_curve
    fig, (ax, ax_h) = plt.subplots(
        2, 1, figsize=(6.0, 4.6), dpi=160,
        gridspec_kw={'height_ratios': [3, 1]}, sharex=True,
    )
    ax.axhline(0, ls='--', color='gray', lw=1, label='идеальная калибровка')
    bins_edges = np.linspace(0, 1, 11)
    series = [
        ('Логистическая регрессия', lr, '#1f77b4', 'o', 2.0, 1.0),
        ('Случайный лес',           rf, '#ff9f1c', 's', 2.0, 0.95),
        ('Градиентный бустинг',     hgb, '#8b0000', '^', 2.0, 0.9),
    ]
    for zi, (name, m, color, marker, lw, alpha) in enumerate(series):
        proba = m.predict_proba(X_test)[:, 1]
        prob_true, prob_pred = calibration_curve(
            y_test, proba, n_bins=10, strategy='uniform',
        )
        residual = prob_true - prob_pred
        ax.plot(prob_pred, residual, marker=marker, lw=lw, markersize=8,
                label=name, color=color, alpha=alpha, zorder=3 + zi,
                markeredgecolor='white', markeredgewidth=0.8)
        ax_h.hist(proba, bins=bins_edges, color=color,
                  alpha=0.5, label=name, edgecolor='white', linewidth=0.5,
                  zorder=3 + zi)
    ax.set_xlim(0, 1)
    ax.set_ylim(-0.12, 0.12)
    ax.set_ylabel('реально − предсказали')
    ax.set_title('Ошибка калибровки по диапазонам вероятности')
    ax.text(0.02, 0.10, 'выше нуля: модель занизила риск',
            fontsize=8.5, color='gray')
    ax.text(0.02, -0.11, 'ниже нуля: модель задирает уверенность',
            fontsize=8.5, color='gray')
    ax.legend(fontsize=8.5, loc='lower right', framealpha=0.9)
    ax_h.set_xlabel('предсказанная вероятность отмены')
    ax_h.set_ylabel('броней (лог.)')
    ax_h.set_yscale('log')
    fig.tight_layout()
    fig.savefig(IMG_DIR / 'calibration.png')
    plt.close(fig)

    # тестовые AUC
    test_auc = {
        'LogReg': roc_auc_score(y_test, lr.predict_proba(X_test)[:, 1]),
        'RandomForest': roc_auc_score(y_test, rf.predict_proba(X_test)[:, 1]),
        'HistGB': roc_auc_score(y_test, hgb.predict_proba(X_test)[:, 1]),
    }
    print('test AUC:', {k: round(v, 4) for k, v in test_auc.items()})


# ----------------- сборка pptx -----------------

def _set_text(ph, text, size=None, bold=None):
    tf = ph.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    if size is not None:
        run.font.size = Pt(size)
    if bold is not None:
        run.font.bold = bold


def _set_bullets(ph, bullets, size=13):
    tf = ph.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, line in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = 0
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)


def _add_image(slide, path, left, top, width, height):
    return slide.shapes.add_picture(str(path), left, top, width=width, height=height)


def _placeholder(slide, idx):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            return ph
    return None


def _textbox(slide, left, top, width, height):
    """Текстовое поле в свободном месте, не привязанное к placeholder."""
    return slide.shapes.add_textbox(left, top, width, height)


def _set_box_text(box, text, size=14, bold=False):
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold


def _set_box_bullets(box, bullets, size=13):
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = 0
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)


def _set_title(slide, text, size=22, width_emu=None):
    """Заголовок слайда. Не использует placeholder (на нём ломалась
    позиция при расширении ширины), вместо этого добавляет свой textbox
    с фиксированной позицией в верхней части слайда и красным цветом СПбГУ.
    """
    # Очищаем шаблонный placeholder, чтобы он не вылезал поверх своим текстом
    ph = _placeholder(slide, 0)
    if ph is not None:
        ph.text_frame.clear()

    # Свой заголовок: широкий textbox по верхней кромке
    w = width_emu if width_emu is not None else 8500000
    box = slide.shapes.add_textbox(Emu(323528), Emu(220000), Emu(w), Emu(550000))
    tf = box.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = SPBGU_RED


def build():
    save_charts()
    pres = Presentation(TEMPLATE)
    layouts = {l.name: l for l in pres.slide_layouts}
    L_title = layouts['TITLE']
    L_work = layouts['1_Рабочий слайд с фотографией']
    L_close = layouts['Закрывающий слайд']

    # стандартные зоны на рабочем слайде (16:9 = 9144000 x 5143500 EMU)
    LEFT_MARGIN = Emu(323528)
    TOP_BELOW_TITLE = Emu(900000)
    # Текстовая колонка слева, картинка справа
    TXT_W = Emu(4200000)
    IMG_LEFT = Emu(4700000)
    IMG_W = Emu(4200000)
    IMG_H = Emu(2800000)
    CONTENT_H = Emu(4000000)
    FULL_W = Emu(8500000)
    SLIDE_W = Emu(9144000)
    SLIDE_H = Emu(5143500)

    # 1. титульный — только название и авторы, без курса
    s = pres.slides.add_slide(L_title)
    _set_text(_placeholder(s, 0),
              'Прогнозирование отмены\nбронирования отеля', size=30, bold=True)
    _set_text(_placeholder(s, 1),
              'Сухоплечев Виталий, Столярова Полина', size=16)

    # вспомогательные функции для рабочего слайда: title через placeholder,
    # подзаголовок и контент - через textbox с явными координатами

    def work_slide(title, subtitle=None, title_width=8500000):
        """Создать рабочий слайд. title_width в EMU: дефолт ~9.3 inch.

        Для длинных заголовков увеличиваем ширину placeholder'а.
        """
        s = pres.slides.add_slide(L_work)
        _set_title(s, title, size=22, width_emu=title_width)
        if subtitle:
            sub = _textbox(s, LEFT_MARGIN, Emu(770000), FULL_W, Emu(380000))
            _set_box_text(sub, subtitle, size=14, bold=False)
        return s

    # 2. задача и зачем это нужно
    s = work_slide('Задача и зачем это нужно',
                   'Предсказать, отменит ли клиент бронь до заезда')
    box = _textbox(s, LEFT_MARGIN, Emu(1200000), FULL_W, CONTENT_H)
    _set_box_bullets(box, [
        '• Что предсказываем: 0 — клиент доехал, 1 — отменил бронь',
        '',
        '• Зачем это отелю.',
        '   Отмены — это потерянные деньги. Если знать риск заранее,',
        '   можно попросить депозит, держать запас по перебронированию',
        '   или прозвонить клиента и подтвердить визит.',
        '',
        '• Чем меряем качество модели.',
        '   ROC-AUC (площадь под ROC-кривой) — основная метрика,',
        '   не зависит от выбранного порога принятия решения.',
        '',
        '   Полнота по классу «отмена» — дополнительная метрика.',
        '   Пропустить отмену хуже, чем ложно её предсказать.',
    ], size=14)

    # 3. данные
    s = work_slide('Данные')
    txt = _textbox(s, LEFT_MARGIN, Emu(1059582), TXT_W, CONTENT_H)
    _set_box_bullets(txt, [
        '• Открытый набор Hotel Booking',
        '   Demand (Antonio и соавт., 2019)',
        '',
        '• 119 390 броней',
        '• 32 признака на бронь',
        '',
        '• Два отеля в Португалии:',
        '   городской и курортный',
        '',
        '• Период: 2015 – 2017 годы',
        '',
        '• Доля отмен в данных — 37 %',
    ], size=14)
    _add_image(s, IMG_DIR / 'balance.png', IMG_LEFT, Emu(1059582), IMG_W, IMG_H)

    # 4. проблемы данных
    s = work_slide('Что почистили в данных',
                   'Без чистки модель ловит шум и утечку целевой переменной')
    box = _textbox(s, LEFT_MARGIN, Emu(1200000), FULL_W, CONTENT_H)
    _set_box_bullets(box, [
        '• Пропуски в колонках:',
        '   страна клиента — 488 пустых → заменили на «неизвестно»;',
        '   агентство — 16 340 пустых → заменили на «нет агентства»;',
        '   компания — 112 593 пустых (94 % колонки!) →',
        '      превратили в бинарный признак «бронь от компании».',
        '',
        '• Цена за ночь:',
        '   встречались значения от −6 € до 5 400 €',
        '   (нормальный диапазон — около 100 €);',
        '   убрали отрицательные и аномально высокие.',
        '',
        '• Дубли: около 32 тысяч одинаковых строк — удалили.',
        '• Брони с нулевым числом гостей — удалили.',
        '',
        '• Утечка целевой переменной:',
        '   колонки «статус брони» и «дата статуса» —',
        '   это прямой ответ, удалили их до обучения.',
    ], size=13)

    # 5. главный сигнал — заголовок длинный, расширяем
    s = work_slide('Главный сигнал: срок до заезда', title_width=8500000)
    txt = _textbox(s, LEFT_MARGIN, Emu(1059582), TXT_W, CONTENT_H)
    _set_box_bullets(txt, [
        '• Чем дольше клиент ждёт',
        '   до заезда, тем чаще отменяет',
        '',
        '• Брони, оформленные за неделю',
        '   до заезда, отменяют примерно',
        '   в 10 % случаев',
        '',
        '• Брони за 4–6 месяцев —',
        '   отменяют почти в половине',
        '',
        '• Это самый сильный',
        '   отдельный признак в данных',
    ], size=13)
    _add_image(s, IMG_DIR / 'lead_time.png', IMG_LEFT, Emu(1059582), IMG_W, IMG_H)

    # 6. депозит и страны
    s = work_slide('Депозит и страна клиента')
    txt = _textbox(s, LEFT_MARGIN, Emu(1059582), TXT_W, CONTENT_H)
    _set_box_bullets(txt, [
        '• Невозвратный депозит —',
        '   почти всегда заканчивается',
        '   отменой. Особенность датасета,',
        '   об этом пишут и авторы статьи.',
        '',
        '• Клиенты из Португалии',
        '   отменяют чаще остальных.',
        '',
        '• Логично: отели находятся',
        '   в Португалии, местные',
        '   бронируют «на всякий случай».',
    ], size=12)
    _add_image(s, IMG_DIR / 'deposit.png', IMG_LEFT, Emu(900000), IMG_W, Emu(1900000))
    _add_image(s, IMG_DIR / 'country.png', IMG_LEFT, Emu(2900000), IMG_W, Emu(1900000))

    # 7. сезонность
    s = work_slide('Сезонность')
    txt = _textbox(s, LEFT_MARGIN, Emu(1059582), TXT_W, CONTENT_H)
    _set_box_bullets(txt, [
        '• Летом доля отмен заметно выше,',
        '   зимой ниже',
        '',
        '• Январь – апрель — самые',
        '   спокойные месяцы',
        '',
        '• Из даты заезда сделали три',
        '   полезных признака:',
        '   – номер месяца',
        '   – время года',
        '   – день недели заезда',
    ], size=13)
    _add_image(s, IMG_DIR / 'season.png', IMG_LEFT, Emu(1059582), IMG_W, IMG_H)

    # 8. подготовка данных и признаков — длинный заголовок, расширяем
    s = work_slide('Подготовка данных и новые признаки',
                   'Код вынесен в модули src/data.py и src/features.py',
                   title_width=8500000)
    box = _textbox(s, LEFT_MARGIN, Emu(1200000), FULL_W, CONTENT_H)
    _set_box_bullets(box, [
        '• Новые признаки, которые мы посчитали:',
        '   общее число ночей и общее число гостей,',
        '   признак «есть ли дети» и «семейная бронь»,',
        '   «сменили номер при заезде» (снижает риск отмены),',
        '   месяц, время года, день недели заезда,',
        '   цена на одного гостя.',
        '',
        '• Категориальные признаки переводим в числа:',
        '   признаки с малым числом значений (тип отеля, питание, канал) —',
        '      двоичное кодирование (one-hot) для линейных моделей,',
        '      порядковое кодирование для деревьев и бустинга;',
        '   страна (около 175 значений) — целевое кодирование:',
        '      значение признака = средняя доля отмен по стране,',
        '      считается отдельно для каждого фолда (без утечки).',
        '',
        '• Разбиение: 80 % обучение / 20 % тест, со стратификацией по цели.',
    ], size=12)

    # 9. сравнение моделей
    s = work_slide('Сравнение моделей')
    txt = _textbox(s, LEFT_MARGIN, Emu(1059582), TXT_W, CONTENT_H)
    _set_box_bullets(txt, [
        'Метрика: ROC-AUC по',
        '5-фолдовой стратифицированной',
        'кросс-валидации.',
        '',
        '• Случайный ответ — 0.50',
        '• Логистическая регрессия — 0.84',
        '• k ближайших соседей — 0.85',
        '• Решающее дерево — 0.88',
        '• Случайный лес — 0.91',
        '• Градиентный бустинг — 0.92',
        '',
        'Разброс по фолдам ≈ 0.002 —',
        'результаты устойчивые.',
    ], size=12)
    _add_image(s, IMG_DIR / 'models.png', IMG_LEFT, Emu(1059582), IMG_W, IMG_H)

    # 10. финальная модель — заголовок длинный, расширяем
    s = work_slide('Финальная модель — градиентный бустинг',
                   title_width=8500000)
    txt = _textbox(s, LEFT_MARGIN, Emu(1059582), TXT_W, CONTENT_H)
    _set_box_bullets(txt, [
        'HistGradientBoosting из',
        'библиотеки scikit-learn',
        '',
        '• Качество на тесте:',
        '   ROC-AUC = 0.917',
        '   F1-мера = 0.77',
        '   Полнота = 0.67',
        '',
        '• Почему выиграл:',
        '   ловит нелинейные связи,',
        '   быстрее случайного леса,',
        '   стабилен по фолдам.',
        '',
        '• Параметры: 300 итераций,',
        '   скорость обучения 0.05,',
        '   глубина дерева 8.',
    ], size=11)
    _add_image(s, IMG_DIR / 'importance.png', IMG_LEFT, Emu(1059582), IMG_W, IMG_H)

    # 11. калибровка — заголовок длинный, расширяем
    s = work_slide('Насколько можно верить вероятностям',
                   title_width=8500000)
    txt = _textbox(s, LEFT_MARGIN, Emu(1059582), TXT_W, CONTENT_H)
    _set_box_bullets(txt, [
        'Идея калибровки: если модель',
        'говорит «70 % шанс», то реально',
        'из таких броней должно',
        'отмениться 70 %.',
        '',
        'На графике справа — насколько',
        'модель промахивается.',
        'Ноль — идеальная калибровка.',
        '',
        '• Логистическая регрессия:',
        '   почти на нуле, промах < 2 %.',
        '',
        '• Случайный лес: уходит ниже',
        '   нуля в средней зоне —',
        '   задирает уверенность до 8 %.',
        '',
        '• Градиентный бустинг —',
        '   похожая картина, до 7 %.',
        '',
        'Снизу видно, что большинство',
        'броней получают вероятность',
        'меньше 0.3 (лог. шкала).',
    ], size=10.5)
    _add_image(s, IMG_DIR / 'calibration.png', IMG_LEFT, Emu(900000), IMG_W, Emu(3800000))

    # 12. демо приложения
    s = work_slide('Демонстрация приложения',
                   'Запуск: run_app.bat — откроется в браузере на localhost:8501')
    box = _textbox(s, LEFT_MARGIN, Emu(1200000), FULL_W, CONTENT_H)
    _set_box_bullets(box, [
        'Сверху страницы — три готовых сценария:',
        '',
        '1. ВЫСОКИЙ РИСК.',
        '   Городской отель, срок 200 дней, невозвратный депозит,',
        '   клиент из Португалии, канал — онлайн-агентство.',
        '   Модель: примерно 99.8 % вероятность отмены.',
        '',
        '2. НИЗКИЙ РИСК.',
        '   Курортный отель, срок 14 дней, прямое бронирование,',
        '   3 спецзапроса, парковка, семья с ребёнком, клиент из Великобритании.',
        '   Модель: примерно 0 % — клиент почти точно доедет.',
        '',
        '3. СРЕДНИЙ РИСК.',
        '   Групповая бронь, срок 90 дней, без спецзапросов.',
        '   Модель: около 57 % — граница риска не острая.',
        '',
        'Дальше — ручная демонстрация: меняем только срок до заезда',
        '(7 → 60 → 200 дней), вероятность отмены растёт почти монотонно.',
    ], size=12)

    # 13. деплой и риски
    s = work_slide('Как развернуть и чем рискуем',
                   'Что нужно для работы в продакшене и слабые места модели',
                   title_width=8500000)
    box = _textbox(s, LEFT_MARGIN, Emu(1200000), FULL_W, CONTENT_H)
    _set_box_bullets(box, [
        '• Развёртывание:',
        '   простой веб-сервис (например, FastAPI) с командой /predict,',
        '   модель загружается из файла в память при старте сервиса,',
        '   время ответа меньше 50 мс на бронь,',
        '   раз в сутки пересчёт по активным бронированиям,',
        '   мониторинг смещения распределений признаков,',
        '   переобучение раз в квартал на скользящем окне.',
        '',
        '• Главные риски:',
        '   1. Перенос на другие отели. Обучались на двух отелях',
        '      в Португалии — на сетевой отель в Германии может',
        '      не перенестись.',
        '   2. Смещение во времени. Данные 2015 – 2017 годов,',
        '      после ковида клиенты ведут себя иначе.',
        '   3. Новая страна. Если в обучении её не было, целевое',
        '      кодирование даёт среднее по миру — точность падает.',
    ], size=12)

    # 14. закрывающий — крупная надпись по центру, красным
    s = pres.slides.add_slide(L_close)
    # очищаем дефолтные placeholder'ы шаблона, чтобы они не светились
    for ph in list(s.placeholders):
        ph.text_frame.clear()
    # центрированная надпись на всю площадь слайда
    box = _textbox(s, Emu(0), Emu(2000000), SLIDE_W, Emu(1200000))
    tf = box.text_frame
    tf.word_wrap = True
    from pptx.enum.text import PP_ALIGN
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = 'Спасибо за внимание'
    run.font.size = Pt(54)
    run.font.bold = True
    run.font.color.rgb = SPBGU_RED

    pres.save(str(OUT_PPTX))
    print('saved:', OUT_PPTX)


if __name__ == '__main__':
    build()
